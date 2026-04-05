import os
import io
from flask import Flask, request, jsonify, send_from_directory, session
from flask_cors import CORS
from werkzeug.security import generate_password_hash, check_password_hash
from google import genai
from dotenv import load_dotenv
from duckduckgo_search import DDGS
from PyPDF2 import PdfReader
import asyncio
from playwright.async_api import async_playwright
import trafilatura
from pptx import Presentation
import fitz  # PyMuPDF
import sqlite3
from datetime import date

# Load environment variables
load_dotenv()

app = Flask(__name__, static_folder='.')
app.secret_key = os.getenv("SECRET_KEY", "dev-secret-12345")
CORS(app)

import time

# Global context for loaded PDF content
document_context = ""
active_file_metadata = None # Stores Gemini File API object or filename

# Initialize Gemini Client with the new SDK
api_key = os.getenv("GEMINI_API_KEY")
if not api_key:
    raise ValueError("GEMINI_API_KEY not found in environment variables")

client = genai.Client(api_key=api_key)

# The user requested Gemini 2.5 flash. 
MODEL_NAME = "gemini-2.5-flash" 

# Setup SQLite for Usage and Auth Tracking
def init_db():
    with sqlite3.connect("metrics.db") as conn:
        cursor = conn.cursor()
        # Users Table
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS users (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                username TEXT UNIQUE,
                password TEXT,
                is_premium INTEGER DEFAULT 0
            )
        ''')
        # Usage Table (Per User)
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS usage (
                user_id INTEGER PRIMARY KEY,
                review_count INTEGER DEFAULT 0,
                last_used_date TEXT
            )
        ''')
        # Conversations Table (Memory)
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS conversations (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                user_id INTEGER,
                role TEXT,
                content TEXT,
                timestamp DATETIME DEFAULT CURRENT_TIMESTAMP
            )
        ''')
        conn.commit()

init_db()

def get_user_review_count(user_id):
    today = str(date.today())
    with sqlite3.connect("metrics.db") as conn:
        cursor = conn.cursor()
        cursor.execute("SELECT review_count, last_used_date FROM usage WHERE user_id = ?", (user_id,))
        row = cursor.fetchone()
        if row:
            count, last_date = row
            if last_date != today:
                cursor.execute("UPDATE usage SET review_count = 0, last_used_date = ? WHERE user_id = ?", (today, user_id))
                return 0
            return count
        else:
            cursor.execute("INSERT INTO usage (user_id, review_count, last_used_date) VALUES (?, 0, ?)", (user_id, today))
            return 0

def increment_review_count(user_id):
    with sqlite3.connect("metrics.db") as conn:
        conn.execute("UPDATE usage SET review_count = review_count + 1 WHERE user_id = ?", (user_id,))
        conn.commit()

def save_message(user_id, role, content):
    with sqlite3.connect("metrics.db") as conn:
        conn.execute("INSERT INTO conversations (user_id, role, content) VALUES (?, ?, ?)", (user_id, role, content))
        conn.commit()

def get_history(user_id, limit=10):
    with sqlite3.connect("metrics.db") as conn:
        cursor = conn.cursor()
        cursor.execute("SELECT role, content FROM conversations WHERE user_id = ? ORDER BY timestamp DESC LIMIT ?", (user_id, limit))
        rows = cursor.fetchall()
        # History needs to be in chronological order for Gemini
        history = [{"role": r[0], "parts": [r[1]]} for r in reversed(rows)]
        return history

# Sequoia Capital Pitch Deck Benchmarks for RAG
SEQUOIA_BENCHMARKS = {
    "Problem": "Define the customer pain point. Explain who experiences this and why current solutions are inadequate.",
    "Solution": "Define your eureka moment and unique value proposition. How do you solve the problem specifically?",
    "Why Now": "Identify market trends, technological shifts, or timing factors that make this the right moment.",
    "Market": "TAM/SAM/SOM analysis. Define your target customer and market size accurately.",
    "Competition": "Detail direct and indirect competitors. Clearly articulate your unique plan to win.",
    "Business Model": "Explain how you intend to make money and sustain growth. Unit economics matter.",
    "Team": "Introduce founders and key members. Highlight relevant expertise and Founder-Market fit.",
    "Financials": "Include key metrics, projections, and clear usage of funds.",
    "Vision": "Describe the long-term impact of your company in 5+ years."
}

async def scrape_single_url(context, url):
    """Helper to scrape a single URL with Playwright and Trafilatura."""
    try:
        page = await context.new_page()
        # Set a 15-second timeout to ensure the search remains fast
        await page.goto(url, wait_until="domcontentloaded", timeout=15000)
        content = await page.content()
        await page.close()
        
        # Trafilatura extracts perfectly clean text (Markdown style)
        extracted = trafilatura.extract(content, include_links=False)
        if extracted:
            return {"success": True, "markdown": extracted}
        return {"success": False, "markdown": ""}
    except Exception as e:
        print(f"Error scraping {url}: {e}")
        return {"success": False, "markdown": ""}

async def _crawl_urls(urls):
    """Asynchronously crawl URLs using Playwright & Trafilatura."""
    crawled_results = []
    try:
        async with async_playwright() as p:
            browser = await p.chromium.launch(headless=True)
            # Create a shared context for speed
            context = await browser.new_context(
                user_agent="Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/122.0.0.0 Safari/537.36"
            )
            
            tasks = [scrape_single_url(context, url) for url in urls]
            results = await asyncio.gather(*tasks, return_exceptions=True)
            
            for r in results:
                if isinstance(r, Exception):
                    crawled_results.append({"success": False, "markdown": ""})
                else:
                    crawled_results.append(r)
            
            await browser.close()
    except Exception as e:
        print(f"Crawl engine error: {e}")
        crawled_results = [{"success": False, "markdown": ""} for _ in urls]
        
    return crawled_results

@app.route('/')
def index():
    return send_from_directory('.', 'index.html')

@app.route('/style.css')
def css():
    return send_from_directory('.', 'style.css')

@app.route('/script.js')
def js():
    return send_from_directory('.', 'script.js')

# --- Authentication Routes ---

@app.route('/api/register', methods=['POST'])
def register():
    data = request.json
    username = data.get('username')
    password = data.get('password')
    if not username or not password:
        return jsonify({"error": "Missing username or password"}), 400
    
    hashed = generate_password_hash(password)
    try:
        with sqlite3.connect("metrics.db") as conn:
            cursor = conn.cursor()
            cursor.execute("INSERT INTO users (username, password) VALUES (?, ?)", (username, hashed))
            conn.commit()
            user_id = cursor.lastrowid
            session['user_id'] = user_id
            session['username'] = username
            return jsonify({"success": True, "user": {"id": user_id, "username": username}})
    except sqlite3.IntegrityError:
        return jsonify({"error": "Username already exists"}), 400

@app.route('/api/login', methods=['POST'])
def login():
    data = request.json
    username = data.get('username')
    password = data.get('password')
    
    with sqlite3.connect("metrics.db") as conn:
        cursor = conn.cursor()
        cursor.execute("SELECT id, password, is_premium FROM users WHERE username = ?", (username,))
        user = cursor.fetchone()
        
        if user and check_password_hash(user[1], password):
            session['user_id'] = user[0]
            session['username'] = username
            return jsonify({"success": True, "user": {"id": user[0], "username": username, "is_premium": user[2]}})
    
    return jsonify({"error": "Invalid username or password"}), 401

@app.route('/api/user_status')
def user_status():
    if 'user_id' in session:
        with sqlite3.connect("metrics.db") as conn:
            cursor = conn.cursor()
            cursor.execute("SELECT is_premium FROM users WHERE id = ?", (session['user_id'],))
            is_premium = cursor.fetchone()[0]
            return jsonify({
                "logged_in": True, 
                "username": session['username'], 
                "is_premium": is_premium
            })
    return jsonify({"logged_in": False})

@app.route('/api/logout', methods=['POST'])
def logout():
    session.clear()
    return jsonify({"success": True})

@app.route('/api/history')
def chat_history_list():
    user_id = session.get('user_id')
    if not user_id:
        return jsonify([])
    
    with sqlite3.connect("metrics.db") as conn:
        cursor = conn.cursor()
        # Fetch the last 10 user questions to act as conversation 'heads'
        cursor.execute("SELECT content, timestamp FROM conversations WHERE user_id = ? AND role = 'user' ORDER BY timestamp DESC LIMIT 10", (user_id,))
        rows = cursor.fetchall()
        
    history = [{"text": r[0], "time": r[1]} for r in rows]
    return jsonify(history)

@app.route('/api/chat', methods=['POST'])
def chat():
    global document_context
    data = request.json
    user_message = data.get('message')
    user_id = session.get('user_id')

    if not user_message:
        return jsonify({"error": "No message provided"}), 400
        
    # Build Chat History if user is logged in
    history = []
    if user_id:
        history = get_history(user_id)
        save_message(user_id, 'user', user_message)

    full_prompt = user_message
    if document_context:
        full_prompt = f"Context from uploaded document:\n{document_context}\n\nUser Question: {user_message}\n\nPlease answer the question based strictly on the provided context if possible."

    try:
        # Use full history for context-aware responses
        contents = history + [{"role": "user", "parts": [full_prompt]}]
        response = client.models.generate_content(
            model=MODEL_NAME,
            contents=contents,
        )

        if user_id:
            save_message(user_id, 'model', response.text)

        return jsonify({
            "response": response.text
        })
    except Exception as e:
        print(f"Error calling Gemini API: {e}")
        return jsonify({"error": str(e)}), 500

@app.route('/api/search', methods=['POST'])
def web_search():
    global document_context
    """AI Search using DuckDuckGo for context, or RAG if PDF is uploaded."""
    data = request.json
    user_message = data.get('message')

    if not user_message:
        return jsonify({"error": "No message provided"}), 400

    try:
        # 1. Use PDF context if available
        user_id = session.get('user_id')
        if document_context:
            full_prompt = f"Context from uploaded document:\n{document_context}\n\nUser Question: {user_message}\n\nPlease answer the question based strictly on the provided context if possible."
            response = client.models.generate_content(
                model=MODEL_NAME,
                contents=full_prompt,
            )
            if user_id:
                save_message(user_id, 'user', user_message)
                save_message(user_id, 'model', response.text)
            return jsonify({
                "response": response.text,
                "sources": [{"title": "Uploaded PDF context", "url": "#pdf"}]
            })

        # 2. Search DuckDuckGo (limit to 2 for speed)
        with DDGS() as ddgs:
            results = list(ddgs.text(user_message, max_results=2))
            
        if not results:
            response = client.models.generate_content(
                model=MODEL_NAME,
                contents=user_message,
            )
            if user_id:
                save_message(user_id, 'user', user_message)
                save_message(user_id, 'model', response.text)
            return jsonify({
                "response": response.text,
                "sources": []
            })
            
        # 3. Crawl the found URLs
        urls_to_crawl = [r['href'] for r in results]
        crawled_data = asyncio.run(_crawl_urls(urls_to_crawl))

        # 4. Build Comprehensive Web Context
        search_context = ""
        for i, (search_dict, crawl_res) in enumerate(zip(results, crawled_data)):
            content = crawl_res["markdown"][:6000] if crawl_res["success"] else search_dict['body']
            search_context += f"\n--- Source {i+1}: {search_dict['title']} ({search_dict['href']}) ---\n{content}\n"

        # 5. Prompt with History & Search Context
        history = []
        user_id = session.get('user_id')
        if user_id:
            history = get_history(user_id)
            save_message(user_id, 'user', user_message)

        full_prompt = f"""You are a helpful AI search assistant. Use the following web search results to answer the user's question accurately.
        Cite your sources using [Source N] notation.

        Web Results:
        {search_context}

        User Question: {user_message}
        """

        contents = history + [{"role": "user", "parts": [full_prompt]}]
        response = client.models.generate_content(
            model=MODEL_NAME,
            contents=contents,
        )

        if user_id:
            save_message(user_id, 'model', response.text)

        return jsonify({
            "response": response.text,
            "sources": [{"title": r['title'], "url": r['href']} for r in results]
        })
    except Exception as e:
        print(f"Error in web search: {e}")
        # Final safety fallback 
        try:
            response = client.models.generate_content(model=MODEL_NAME, contents=user_message)
            return jsonify({"response": response.text, "sources": []})
        except:
            return jsonify({"error": str(e)}), 500

@app.route('/api/upload', methods=['POST'])
def upload_file():
    global document_context, active_file_metadata
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400

    if file and (file.filename.endswith('.pdf') or file.filename.endswith('.pptx') or file.filename.endswith('.ppt')):
        try:
            file_extension = file.filename.split('.')[-1].lower()
            file_content = file.read()
            
            if file_extension == 'pdf':
                # Save to a temporary file for Gemini File API and PyMuPDF
                temp_filename = f"temp_{file.filename}"
                with open(temp_filename, "wb") as f:
                    f.write(file_content)
                
                # 1. Advanced Text Extraction with PyMuPDF (fitz)
                doc = fitz.open(temp_filename)
                text = ""
                for i, page in enumerate(doc):
                    page_text = page.get_text("text")
                    text += f"\n--- SLIDE {i+1} ---\n{page_text}\n"
                document_context = text
                
                # 2. Upload to Gemini File API for Multimodal Vision
                print(f"Uploading {file.filename} to Gemini File API...")
                uploaded_file = client.files.upload(file=temp_filename)
                
                # Cleanup temp file
                if os.path.exists(temp_filename):
                    os.remove(temp_filename)
                    
            else:
                # Handle PPTX/PPT with Speaker Notes
                pptx_stream = io.BytesIO(file_content)
                prs = Presentation(pptx_stream)
                text = ""
                for slide in prs.slides:
                    # Slide Content
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                    # Speaker Notes
                    if slide.has_notes_slide:
                        notes = slide.notes_slide.notes_text_frame.text
                        text += f"\n[Speaker Notes]: {notes}\n"
                
                document_context = text
                active_file_metadata = {"type": "pptx", "filename": file.filename}
            
            return jsonify({"success": True, "filename": file.filename, "text_length": len(document_context)})
        except Exception as e:
            print(f"Error processing file: {e}")
            return jsonify({"error": str(e)}), 500

    return jsonify({"error": "File type not supported. Please upload a PDF or PowerPoint file."}), 400

@app.route('/api/review', methods=['POST'])
def review_deck():
    user_id = session.get('user_id')
    if not user_id:
        return jsonify({"error": "AUTH_REQUIRED", "message": "Please login to use the Pitch Deck Reviewer."}), 401

    # Check Premium status
    with sqlite3.connect("metrics.db") as conn:
        cursor = conn.cursor()
        cursor.execute("SELECT is_premium FROM users WHERE id = ?", (user_id,))
        is_premium = cursor.fetchone()[0]

    if not is_premium:
        count = get_user_review_count(user_id)
        if count >= 3:
            return jsonify({
                "error": "LIMIT_REACHED", 
                "message": "You have reached your free trial limit of 3 reviews. Please upgrade to Pro for unlimited access."
            }), 429
        increment_review_count(user_id)

    global document_context, active_file_metadata
    if not document_context and not active_file_metadata:
        return jsonify({"error": "No document uploaded for review"}), 400
        # Prepare RAG context from Sequoia Benchmarks
    benchmarks_context = "\n".join([f"- {k}: {v}" for k, v in SEQUOIA_BENCHMARKS.items()])

    prompt = f"""
You are a Senior Partner at Sequoia Capital. You are performing a high-stakes "Deep Audit" of a pitch deck. 
Your evaluation must be structured into 5 distinct "Engines" to provide a comprehensive investment thesis.

1. **Review Engine (Structured Output)**: A professional, concise executive summary and a clear investment verdict.
2. **Scoring Engine**: Numeric evaluation against our 10-point internal roadmap.
3. **Contrarian Engine**: The "Devil's Advocate". Identify exactly why this company will fail. Be brutally honest.
4. **Improvement Engine**: Actionable, high-impact fixes for the deck and the business model.
5. **Pitch Engine**: A single, punchy "One-line Pitch" that captures the essence of the value proposition.

SEQUOIA BENCHMARKS (RAG Context):
{benchmarks_context}

**Instructions for Output:**
1. You MUST output your response in strict JSON format.
2. For each scoring category, provide a numeric score from 0.0 to 5.0 (as a float), and a 'star_rating' string exactly matching the format `★★★☆☆` (e.g. 3.0 = `★★★☆☆`, 4.5 = `★★★★☆`).
3. Follow this exact JSON schema:

{{
  "one_line_pitch": "string",
  "executive_summary": "string",
  "verdict": "string",
  "scoring_engine": {{
    "problem_clarity": {{"numeric_score": 0.0, "star_rating": "string", "justification": "string"}},
    "solution_strength": {{"numeric_score": 0.0, "star_rating": "string", "justification": "string"}},
    "market_opportunity": {{"numeric_score": 0.0, "star_rating": "string", "justification": "string"}},
    "business_model": {{"numeric_score": 0.0, "star_rating": "string", "justification": "string"}},
    "competitive_advantage": {{"numeric_score": 0.0, "star_rating": "string", "justification": "string"}},
    "team_strength": {{"numeric_score": 0.0, "star_rating": "string", "justification": "string"}},
    "financial_clarity": {{"numeric_score": 0.0, "star_rating": "string", "justification": "string"}},
    "storytelling": {{"numeric_score": 0.0, "star_rating": "string", "justification": "string"}},
    "innovation": {{"numeric_score": 0.0, "star_rating": "string", "justification": "string"}},
    "overall_quality": {{"numeric_score": 0.0, "star_rating": "string", "justification": "string"}}
  }},
  "contrarian_engine": {{
    "failure_modes": ["string", "string"],
    "devils_advocate_thesis": "string"
  }},
  "improvement_engine": {{
    "critical_fixes": ["string"],
    "strategic_next_steps": ["string"]
  }},
  "graph_data": {{
    "radar_chart": [0.0, 0.0, 0.0, 0.0, 0.0, 0.0, 0.0, 0.0, 0.0, 0.0], 
    "bar_graph": [0.0, 0.0, 0.0, 0.0, 0.0, 0.0, 0.0, 0.0, 0.0, 0.0]
  }}
}}

Note for graph_data: Make sure the arrays have exactly 10 floats in this specific order: problem_clarity, solution_strength, market_opportunity, business_model, competitive_advantage, team_strength, financial_clarity, storytelling, innovation, overall_quality.
"""

    try:
        contents = [prompt]
        if active_file_metadata and active_file_metadata.get("type") == "pdf":
            contents.append(active_file_metadata["file_obj"])
        else:
            contents.append(f"Pitch Deck Content:\n{document_context[:15000]}")

        from google.genai import types
        response = client.models.generate_content(
            model=MODEL_NAME,
            contents=contents,
            config=types.GenerateContentConfig(
                response_mime_type="application/json"
            )
        )
        
        # Cleanup Gemini file if it exists
        if active_file_metadata and active_file_metadata.get("type") == "pdf":
            try:
                client.files.delete(name=active_file_metadata["file_obj"].name)
                print(f"Deleted Gemini file: {active_file_metadata['file_obj'].name}")
            except Exception as e:
                print(f"Failed to delete Gemini file: {e}")

        if user_id:
            save_message(user_id, 'user', f"Review Pitch Deck: {active_file_metadata.get('filename', 'Unknown File')}")
            save_message(user_id, 'model', "Deck successfully audited. Insights generated.")

        return jsonify({
            "response": response.text
        })
    except Exception as e:
        print(f"Error in pitch deck review: {e}")
        return jsonify({"error": str(e)}), 500

@app.route('/api/clear_context', methods=['POST'])
def clear_context():
    global document_context
    document_context = ""
    return jsonify({"success": True})

if __name__ == '__main__':
    # Running on 0.0.0.0 to ensure accessibility
    app.run(debug=True, port=5000, host='0.0.0.0')
